import os
import re
import subprocess
from typing import List, Tuple, TypedDict

from PIL import Image
from pptx import Presentation

__all__ = ["build_pptx", "slidev2pptx"]


class Page(TypedDict):
    image_path: str
    slide_index: int
    click_index: int
    notes: str


def is_slide(line: str) -> bool:
    """Check if a line corresponds to a slide."""

    return line.startswith("![")


def extract_image_path(line: str) -> str:
    """Extract the image path from a line."""

    # Regular expression pattern to match the file path
    pattern = r"\((./.*?.png)\)"
    return re.search(pattern, line).group(1)


def build_pages(notes: str) -> List[Page]:
    """Extract metadata for each slide page from the presenter notes."""

    lines = notes.split("\n")
    pages = []
    last_line_is_slide = False
    current_notes = ""
    for i, line in enumerate(lines):
        if last_line_is_slide:
            # If the last line was a slide, the current line should be empty
            assert line == "", f"Expected empty line after slide, got: {line}"
            last_line_is_slide = False
            continue

        if not is_slide(line):
            # If the current line is not a slide, it is a part of presenter notes.
            current_notes += line + "\n"

        if is_slide(line) or i == len(lines) - 1:
            # Assign the current notes to the latest pages with the same slide_index.
            if current_notes != "":
                assert len(pages) > 0, "Notes without a slide"
                latest_slide_index = pages[-1]["slide_index"]
                latest_pages = [
                    page for page in pages if page["slide_index"] == latest_slide_index
                ]
                notes_per_click = current_notes.split("[click]")
                for j, notes in enumerate(notes_per_click):
                    latest_pages[j]["notes"] = notes.lstrip()
                current_notes = ""

        if is_slide(line):
            # Create a new slide object.
            # image_path example: "./slides-export\001-01.png"
            image_path = extract_image_path(line)
            image_filename = os.path.basename(image_path)
            slide_index = int(image_filename.split("-")[0])
            click_index = int(image_filename.split("-")[1].split(".png")[0])
            page = Page(
                image_path=image_path,
                slide_index=slide_index,
                click_index=click_index,
                notes="",
            )
            pages.append(page)
            last_line_is_slide = True
    return pages


def get_slides_size(image_directory: str) -> Tuple[int, int]:
    """Get the size of the slides in the image directory."""

    files = os.listdir(image_directory)
    assert len(files) > 0, "Image directory empty"
    first_file = files[0]
    assert first_file.endswith(".png"), f"Expected .png file, got: {first_file}"

    image_path = os.path.join(image_directory, first_file)
    with Image.open(image_path) as img:
        width, height = img.size
    return width, height


def build_pptx(notes_path: str, image_directory: str, output_path: str) -> None:
    """
    Build a PPTX presentation from the exported Slidev notes and images.

    Parameters
    ----------
    notes_path : str
        The path to the exported Slidev presenter notes.
    image_directory : str
        The directory containing the exported Slidev slides as PNG images.
    output_path : str
        The path to save the PPTX presentation.
    """

    # Parse the markdown file to extract the image path and presenter notes for each slide.
    with open(notes_path, "r", encoding="utf-8") as f:
        notes = f.read()
    pages = build_pages(notes)

    # Get the aspect ratio of the slides.
    width, height = get_slides_size(image_directory)
    aspect_ratio = width / height

    # Create a PPTX file with each image as a slide and add notes.
    presentation = Presentation()
    presentation.slide_width = int(presentation.slide_height * aspect_ratio)
    for page in pages:
        # Create empty slide.
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        # Add image.
        image_path = os.path.join(os.path.dirname(image_directory), page["image_path"])
        slide.shapes.add_picture(image_path, 0, 0, height=presentation.slide_height)

        # Add notes.
        slide.notes_slide.notes_text_frame.text = page["notes"]
    presentation.save(output_path)


def run_command(command: str, directory: str = "./") -> Tuple[str, str]:
    """
    Run a command in a given directory and print the output.

    Parameters
    ----------
    command : str
        The command to run.
    directory : str
        The directory in which to run the command.
    """

    process = subprocess.Popen(
        command,
        shell=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        cwd=directory,
    )

    output, error = process.communicate()
    output = output.decode()
    error = error.decode()
    return output, error


def slidev2pptx(slidev_path: str, output_path: str, scale: int = 2) -> None:
    """
    Convert a Slidev slide deck to a PPTX slide deck.

    Parameters
    ----------
    slidev_path : str
        The path to the Slidev slide deck repository.
    output_path : str
        The path to save the PPTX slide deck.
    """

    # Export slides as PNG images for each slide and a markdown file of the notes.
    command = f"pnpm run export --format md --with-clicks --scale {scale}"
    output, error = run_command(command, slidev_path)

    print("---------------------")
    print("Slidev export output:")
    print(output)
    print("---------------------")
    if len(error) > 0:
        print("Slidev export error:")
        print(error)
        print("---------------------")

    notes_path = os.path.join(slidev_path, "slides-export.md")
    image_directory = os.path.join(slidev_path, "slides-export")
    build_pptx(notes_path, image_directory, output_path)
